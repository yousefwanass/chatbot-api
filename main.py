import os
import time
import re
import uuid
import threading
from io import BytesIO
from typing import List, Dict, Optional, Any

import numpy as np
import faiss
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Query
from pydantic import BaseModel
from dotenv import load_dotenv

from PyPDF2 import PdfReader
from sentence_transformers import SentenceTransformer
import google.generativeai as genai
from pptx import Presentation

load_dotenv()

# Optional libs
try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    import pytesseract
    from PIL import Image, ImageOps
except Exception:
    pytesseract = None
    Image = None
    ImageOps = None


# ---------------- CONFIG ----------------
APP_NAME = "Lecture Chatbot API"
APP_VERSION = "1.2.0"

MODEL_NAME = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
EMBED_MODEL_NAME = os.getenv("EMBED_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
TOP_K = int(os.getenv("TOP_K", "4"))
MAX_CHUNK_WORDS = int(os.getenv("MAX_CHUNK_WORDS", "500"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "60"))
MAX_HISTORY_TURNS = int(os.getenv("MAX_HISTORY_TURNS", "20"))

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    raise RuntimeError("Missing GEMINI_API_KEY env var. Set it before running.")
genai.configure(api_key=GEMINI_API_KEY)


# ---------------- APP ----------------
app = FastAPI(title=APP_NAME, version=APP_VERSION)


# ---------------- API CONTRACT ----------------
def _meta() -> Dict[str, Any]:
    return {
        "service": APP_NAME,
        "version": APP_VERSION,
        "time": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    }


def ok(data: Any, *, legacy: bool) -> Any:
    return data if legacy else {"success": True, "data": data, "meta": _meta()}


# ---------------- SESSION STORE ----------------
# sessionId -> { index, metas, lecture_name, chat_history, last_mcqs, created_at }
SESSIONS: Dict[str, Dict[str, Any]] = {}
SESSIONS_LOCK = threading.Lock()


# ---------------- REQUEST/RESPONSE SCHEMAS ----------------
class UploadResponse(BaseModel):
    sessionId: str
    lectureName: str
    chunks: int


class ChatRequest(BaseModel):
    sessionId: str
    message: str
    includeSources: bool = True


class ChatResponse(BaseModel):
    type: str  # "chat" | "mcq" | "answers"
    reply: str
    sources: Optional[str] = None


# ---------------- EXTRACTION HELPERS ----------------
def extract_text_from_pdf_bytes(file_bytes: bytes) -> str:
    try:
        reader = PdfReader(BytesIO(file_bytes))
        parts = []
        for page in reader.pages:
            txt = page.extract_text()
            if txt:
                parts.append(txt)
        return "\n".join(parts).strip()
    except Exception:
        return ""


def extract_text_from_pptx_bytes(file_bytes: bytes) -> str:
    try:
        prs = Presentation(BytesIO(file_bytes))
        full_text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        txt = shape.text.strip()
                        if txt:
                            full_text.append(txt)
                except Exception:
                    pass

                try:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs).strip()
                            if line:
                                full_text.append(line)
                except Exception:
                    pass

                try:
                    if getattr(shape, "has_table", False) and shape.table is not None:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                cell_text = (cell.text or "").strip()
                                if cell_text:
                                    full_text.append(cell_text)
                except Exception:
                    pass

            try:
                if getattr(slide, "has_notes_slide", False) and slide.notes_slide is not None:
                    notes_tf = slide.notes_slide.notes_text_frame
                    if notes_tf is not None:
                        notes = (notes_tf.text or "").strip()
                        if notes:
                            full_text.append(notes)
            except Exception:
                pass

        return "\n".join(full_text).strip()
    except Exception:
        return ""


def extract_text_from_docx_bytes(file_bytes: bytes) -> str:
    if DocxDocument is None:
        return ""
    try:
        doc = DocxDocument(BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        paragraphs.append(cell.text.strip())

        return "\n".join(paragraphs).strip()
    except Exception:
        return ""


def extract_text_from_image_ocr_bytes(file_bytes: bytes) -> str:
    if pytesseract is None or Image is None or ImageOps is None:
        return ""
    try:
        img = Image.open(BytesIO(file_bytes)).convert("RGB")
        img = ImageOps.grayscale(img)
        img = ImageOps.autocontrast(img)
        text = pytesseract.image_to_string(img, lang="eng+ara")
        return (text or "").strip()
    except Exception:
        return ""


def extract_text_by_extension(file_bytes: bytes, ext: str, filename: str) -> str:
    if ext == ".pdf":
        return extract_text_from_pdf_bytes(file_bytes)

    if ext == ".pptx":
        text = extract_text_from_pptx_bytes(file_bytes)
        if not text:
            raise HTTPException(status_code=400, detail=f"{filename}: could not read PPTX.")
        return text

    if ext == ".docx":
        if DocxDocument is None:
            raise HTTPException(status_code=400, detail="python-docx not installed; can't read DOCX.")
        return extract_text_from_docx_bytes(file_bytes)

    if ext in {".jpg", ".jpeg", ".png"}:
        return extract_text_from_image_ocr_bytes(file_bytes) if pytesseract else ""

    return ""


# ---------------- TEXT / RAG HELPERS ----------------
def chunk_text(text: str, max_len: int = MAX_CHUNK_WORDS, overlap: int = CHUNK_OVERLAP) -> List[str]:
    if not text:
        return []

    words = text.split()
    chunks = []
    step = max(1, max_len - overlap)
    i = 0

    while i < len(words):
        chunk = words[i:i + max_len]
        if chunk:
            chunks.append(" ".join(chunk))
        i += step

    return chunks


_embed_model: Optional[SentenceTransformer] = None
_embed_lock = threading.Lock()


def get_embed_model() -> SentenceTransformer:
    global _embed_model
    with _embed_lock:
        if _embed_model is None:
            _embed_model = SentenceTransformer(EMBED_MODEL_NAME)
        return _embed_model


def embed_chunks(chunks: List[str]) -> np.ndarray:
    model = get_embed_model()
    return model.encode(chunks, convert_to_numpy=True, show_progress_bar=False)


def build_faiss_index(embeddings: np.ndarray) -> faiss.IndexFlatIP:
    index = faiss.IndexFlatIP(embeddings.shape[1])
    faiss.normalize_L2(embeddings)
    index.add(embeddings)
    return index


def retrieve_docs(query: str, index: faiss.IndexFlatIP, metas: List[Dict[str, Any]], top_k: int = TOP_K) -> List[Dict[str, Any]]:
    model = get_embed_model()
    q_emb = model.encode([query], convert_to_numpy=True)
    faiss.normalize_L2(q_emb)
    distances, indices = index.search(q_emb, top_k)

    docs = []
    for score, idx in zip(distances[0], indices[0]):
        if 0 <= idx < len(metas):
            docs.append({"score": float(score), "text": metas[idx]["text"]})
    return docs


def expand_short_query(q: str) -> str:
    q = (q or "").strip()
    if len(q.split()) <= 2 and not any(w in q.lower() for w in ["define", "what", "explain"]):
        return f"Explain {q}."
    return q


def clean_model_text(txt: str) -> str:
    if not txt:
        return ""
    txt = re.sub(r"<\|.*?\|>", " ", txt)
    txt = txt.replace("begin_of_sentence", " ").replace("end_of_sentence", " ")
    txt = re.sub(r"([^\w\s])\1{6,}", r"\1", txt)
    txt = "\n".join(re.sub(r"[ \t]{2,}", " ", line).strip() for line in txt.splitlines())
    txt = re.sub(r"(#+\s)", r"\n\n\1", txt)
    txt = re.sub(r"(\|.*\|)", r"\n\1\n", txt)
    return txt.strip()


def format_sources(docs: List[Dict[str, Any]], max_chars: int = 240) -> str:
    if not docs:
        return ""
    lines = []
    for i, d in enumerate(docs, 1):
        snippet = (d.get("text") or "").replace("\n", " ").strip()
        if len(snippet) > max_chars:
            snippet = snippet[:max_chars] + "..."
        lines.append(f"- Source {i} (score={d.get('score', 0):.2f}): {snippet}")
    return "\n".join(lines)


# ---------------- PROMPTS ----------------
def make_prompt(question: str, docs: List[Dict[str, Any]], lecture_name: str) -> str:
    context = "\n\n".join(f"Context [{i+1}]: {d['text']}" for i, d in enumerate(docs))
    return f"""
You are a friendly, intelligent university assistant.
You are answering based on the lecture file **{lecture_name}**.
Only use the provided context to answer questions.
If the context doesn’t contain enough information, say politely:
"I'm sorry, I couldn't find that information in your lecture."

{context}

Question: {question}

Answer clearly, concisely, and with proper markdown formatting.
""".strip()


def make_mcq_prompt(docs: List[Dict[str, Any]], lecture_name: str, num_questions: int = 5) -> str:
    context = "\n\n".join(f"Context [{i+1}]: {d['text']}" for i, d in enumerate(docs))
    return f"""
You are a skilled university tutor creating multiple-choice questions.
Based only on the lecture file **{lecture_name}**, generate **{num_questions}** high-quality MCQs.

Each question must:
- Focus on conceptual understanding
- Have exactly 4 options (A–D)
- Put each option on its OWN line
- Separate each question with two blank lines

Follow this EXACT format:
Q1. <question>
A) <option A>
B) <option B>
C) <option C>
D) <option D

Q2. <question>
A) <option A>
B) <option B>
C) <option C>
D) <option D

Important:
- NEVER put options A–D on the same line.
- Each must start with its label on a new line.
- Do not bold, italicize, or restyle anything.

{context}
""".strip()


def make_answer_prompt(previous_mcqs: str) -> str:
    return f"""
You are a precise evaluator.
The user previously generated multiple-choice questions.
Now, ONLY provide the correct answers for those questions.

Rules:
- Do NOT repeat questions or options.
- ONLY output answers like this (each on a new line):
Q1: B
Q2: D
Q3: C
Q4: A
Q5: B

Here are the questions:
{previous_mcqs}
""".strip()


# ---------------- GEMINI ----------------
def call_gemini(prompt: str) -> str:
    model = genai.GenerativeModel(MODEL_NAME)

    for _ in range(2):
        try:
            response = model.generate_content(prompt)

            if hasattr(response, "text") and response.text:
                return response.text

            candidates = getattr(response, "candidates", None)
            if candidates:
                parts = getattr(candidates[0].content, "parts", None)
                if parts and len(parts) > 0 and hasattr(parts[0], "text"):
                    return parts[0].text
        except Exception:
            time.sleep(0.15)

    return "⚠️ Sorry, I couldn't process your request right now."


# ---------------- SESSION HELPERS ----------------
def get_session(session_id: str) -> Dict[str, Any]:
    with SESSIONS_LOCK:
        sess = SESSIONS.get(session_id)
    if not sess:
        raise HTTPException(status_code=404, detail="Invalid sessionId. Upload files first.")
    return sess


def append_history(sess: Dict[str, Any], role: str, content: str) -> None:
    sess["chat_history"].append({"role": role, "content": content})

    max_items = MAX_HISTORY_TURNS * 2
    if len(sess["chat_history"]) > max_items:
        sess["chat_history"] = sess["chat_history"][-max_items:]


# ---------------- ENDPOINTS ----------------
@app.get("/health")
def health(format: str = Query("legacy", pattern="^(legacy|standard)$")):
    legacy = (format != "standard")
    payload = {
        "status": "ok",
        "service": APP_NAME,
        "version": APP_VERSION,
        "sessionsLoaded": len(SESSIONS),
        "ocrAvailable": pytesseract is not None,
        "docxAvailable": DocxDocument is not None,
    }
    return ok(payload, legacy=legacy)


@app.post("/ai/upload", response_model=UploadResponse)
async def upload(
    files: List[UploadFile] = File(...),
    sessionId: Optional[str] = Form(default=None),
):
    sid = sessionId or str(uuid.uuid4())
    allowed = {".pdf", ".pptx", ".docx", ".jpg", ".jpeg", ".png"}

    all_chunks: List[str] = []
    lecture_names: List[str] = []

    for f in files:
        name = f.filename or "file"
        lecture_names.append(name)
        ext = os.path.splitext(name.lower())[1]

        if ext not in allowed:
            raise HTTPException(status_code=400, detail=f"Unsupported file: {name}")

        file_bytes = await f.read()
        text = extract_text_by_extension(file_bytes, ext, name)

        if text:
            all_chunks.extend(chunk_text(text))

    if not all_chunks:
        raise HTTPException(status_code=400, detail="No text extracted. Check files or OCR/docx dependencies.")

    embeddings = np.array(embed_chunks(all_chunks))
    if embeddings.ndim != 2 or embeddings.size == 0:
        raise HTTPException(status_code=500, detail="Embedding generation failed.")

    index = build_faiss_index(embeddings)
    metas = [{"text": c} for c in all_chunks]
    lecture_name = ", ".join(lecture_names)

    with SESSIONS_LOCK:
        SESSIONS[sid] = {
            "index": index,
            "metas": metas,
            "lecture_name": lecture_name,
            "chat_history": [],
            "last_mcqs": "",
            "created_at": time.time(),
        }

    return UploadResponse(sessionId=sid, lectureName=lecture_name, chunks=len(all_chunks))


@app.post("/ai/chat", response_model=ChatResponse)
def chat(req: ChatRequest):
    sess = get_session(req.sessionId)

    message = (req.message or "").strip()
    if not message:
        raise HTTPException(status_code=400, detail="message is required")

    msg_low = message.lower()

    wants_answers = any(w in msg_low for w in ["answer", "answers", "solve", "solutions", "correct answers"])
    wants_mcq = any(w in msg_low for w in ["mcq", "quiz", "questions", "question", "test"])

    # 1) Answers mode
    if wants_answers and (sess.get("last_mcqs") or "").strip():
        prompt = make_answer_prompt(sess["last_mcqs"])
        raw = call_gemini(prompt)

        answers = clean_model_text(raw)
        answers = re.sub(r"(Q\d+:)", r"\n\1", answers).strip()

        append_history(sess, "user", message)
        append_history(sess, "assistant", answers)

        return ChatResponse(type="answers", reply=answers, sources=None)

    # 2) MCQ mode
    if wants_mcq:
        m = re.search(r"(\d+)\s*(mcq|question|questions|quiz|test)", msg_low)
        n = int(m.group(1)) if m else 5
        n = max(1, min(20, n))

        docs = retrieve_docs(expand_short_query(message), sess["index"], sess["metas"], top_k=TOP_K)
        prompt = make_mcq_prompt(docs, sess["lecture_name"], num_questions=n)
        raw = call_gemini(prompt)

        mcqs = clean_model_text(raw)
        mcqs = re.sub(r"([A-D]\))", r"\n\1", mcqs)
        mcqs = re.sub(r"(D\).*?)(?=Q\d+\.)", r"\1\n\n", mcqs, flags=re.DOTALL)

        sess["last_mcqs"] = mcqs

        append_history(sess, "user", message)
        append_history(sess, "assistant", mcqs)

        return ChatResponse(type="mcq", reply=mcqs, sources=None)

    # 3) Normal chat mode
    expanded = expand_short_query(message)
    docs = retrieve_docs(expanded, sess["index"], sess["metas"], top_k=TOP_K)

    history_text = "\n".join(
        f"{turn['role']}: {turn['content']}"
        for turn in sess["chat_history"][-(MAX_HISTORY_TURNS * 2):]
    )

    q_text = expanded + (("\n\nPrevious chat:\n" + history_text) if history_text else "")
    prompt = make_prompt(q_text, docs, sess["lecture_name"])

    raw = call_gemini(prompt)
    reply = clean_model_text(raw)

    append_history(sess, "user", message)
    append_history(sess, "assistant", reply)

    sources = format_sources(docs) if req.includeSources else None
    return ChatResponse(type="chat", reply=reply, sources=sources)